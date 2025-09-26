# python/make_deck.py
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parents[1]
DATA = ROOT / "data"

def add_title_slide(prs, title, subtitle=None):
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle

def add_picture_slide(prs, title, img_path):
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    left = Inches(0.5)
    top = Inches(1.3)
    height = Inches(5.2)  # keep aspect ratio by setting only one dimension
    slide.shapes.add_picture(str(img_path), left, top, height=height)

def add_text_slide(prs, title, text):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(text.splitlines()):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = line
        p.level = 0
        p.font.size = Pt(14)
    # center-ish
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.LEFT

def main():
    prs = Presentation()
    add_title_slide(prs, "FP&A AI Dashboard – Executive Pack", "Automated Variance • Forecast • AI Summary")

    # Executive summary
    exec_md = (DATA / "exec_summary.md").read_text(encoding="utf-8") if (DATA / "exec_summary.md").exists() else "Summary unavailable."
    add_text_slide(prs, "Executive Summary", exec_md)

    # Visuals (ensure they exist)
    charts = [
        ("Revenue vs Budget (Trend)", DATA / "viz_trend.png"),
        ("Variance % by Department", DATA / "viz_dept_variance.png"),
        ("Forecast by Department (6 mo.)", DATA / "viz_forecast_dept.png"),
    ]
    for title, path in charts:
        if path.exists():
            add_picture_slide(prs, title, path)

    out = DATA / "fpna_onepager.pptx"
    prs.save(out)
    print("✅ Deck created:", out)

if __name__ == "__main__":
    main()
